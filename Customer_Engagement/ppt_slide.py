from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # Widescreen 16:9
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]

# Theme Colors
NAVY = RGBColor(16, 44, 87)
BLUE = RGBColor(0, 114, 206)
RED = RGBColor(230, 57, 70)
GREEN = RGBColor(42, 157, 143)
GRAY = RGBColor(100, 100, 100)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)


def add_header(slide, title_text, subtitle_text=""):
    """Helper to create consistent standard headers across slides."""
    tb = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide1 = prs.slides.add_slide(blank_slide_layout)

# Background accent box
rect = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
)
rect.fill.solid()
rect.fill.fore_color.rgb = NAVY
rect.line.fill.background()

tb = slide1.shapes.add_textbox(
    Inches(1.5), Inches(2.5), Inches(10.333), Inches(3.0)
)
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "European Bank"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.text = "Customer Engagement & Retention Analytics"
p2.font.size = Pt(28)
p2.font.color.rgb = BLUE

p3 = tf.add_paragraph()
p3.text = "\nDashboard Executive Summary & Profile Analysis"
p3.font.size = Pt(16)
p3.font.color.rgb = WHITE

# ==========================================
# SLIDE 2: Global Filters & Executive KPIs
# ==========================================
slide2 = prs.slides.add_slide(blank_slide_layout)
add_header(
    slide2,
    "Executive Overview & Applied Dashboard Filters",
    "High-level metrics and current active parameters",
)

# Left Box: Global Filters
filter_box = slide2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8),
    Inches(1.8),
    Inches(3.8),
    Inches(5.0),
)
filter_box.fill.solid()
filter_box.fill.fore_color.rgb = LIGHT_BG
filter_box.line.color.rgb = BLUE

tf_f = filter_box.text_frame
tf_f.word_wrap = True
p = tf_f.paragraphs[0]
p.text = "🔍 Global Filters Applied"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY

items = [
    ("\nGeographies Selected:", True),
    (" • France\n • Spain\n • Germany", False),
    ("\nHigh Balance Cutoff:", True),
    (" • $127,644.18", False),
    ("\nSource Dataset:", True),
    (" • European_Bank.csv", False),
]
for text, is_bold in items:
    p = tf_f.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = is_bold
    p.font.color.rgb = NAVY if is_bold else GRAY

# Right Cards: KPIs
kpis = [
    ("Total Customers", "10,000", NAVY),
    ("Overall Churn Rate", "20.4%", RED),
    ("Active / Inactive Churn", "14.3% / 26.9%", BLUE),
    ("At-Risk Premium Customers", "1,247", RED),
]

for i, (title, val, color) in enumerate(kpis):
    x = Inches(5.0 + (i % 2) * 3.8)
    y = Inches(1.8 + (i // 2) * 2.5)

    card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(2.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color

    tf_c = card.text_frame
    tf_c.word_wrap = True

    p1 = tf_c.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = GRAY

    p2 = tf_c.add_paragraph()
    p2.text = f"\n{val}"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = color

# ==========================================
# SLIDE 3: Engagement Profiles Analysis
# ==========================================
slide3 = prs.slides.add_slide(blank_slide_layout)
add_header(
    slide3,
    "🎯 Engagement Classification & Risk Analysis",
    "Breakdown of customer share and churn risk across engagement profiles",
)

# Table for Data
rows, cols = 6, 3
left, top, width, height = (
    Inches(0.8),
    Inches(1.8),
    Inches(11.7),
    Inches(4.8),
)
table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(4.5)
table.columns[1].width = Inches(3.6)
table.columns[2].width = Inches(3.6)

headers = ["Engagement Profile", "% Customer Share", "Churn Rate (%)"]
data = [
    ("Active Engaged", "25.9%", "~9.8% (Lowest Risk)"),
    ("Active Low-Product", "25.6%", "~18.9%"),
    ("Moderate Engagement", "19.1%", "~14.5%"),
    ("Inactive Disengaged", "16.9%", "~38.2% (Highest Risk)"),
    ("Inactive High-Balance", "12.5%", "~30.5% (Critical Value Risk)"),
]

# Header Row Formatting
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = WHITE

# Data Rows Formatting
for row_idx, row_data in enumerate(data, start=1):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = (
            LIGHT_BG if row_idx % 2 == 0 else WHITE
        )  # Alternating colors
        for p in cell.text_frame.paragraphs:
            if col_idx > 0:
                p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(14)
                run.font.color.rgb = NAVY

# Save presentation
prs.save("European_Bank_Retention_Analytics.pptx")
print(
    "Presentation successfully generated and saved as 'European_Bank_Retention_Analytics.pptx'"
)